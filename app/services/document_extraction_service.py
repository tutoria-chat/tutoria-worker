"""
Document Text Extraction Service

Extracts text from files for RAG performance optimization.
Stores extracted text in database to avoid re-extraction on every chat message.

Supports:
- Local extraction using pdfplumber (fast, free - PDFs only)
- AI-powered extraction with provider fallback chain (any file type)
- Automatic summarization for long content

Provider chain is configured via AIModels table (UseForFileExtraction=True),
falling back to config settings if no DB models are flagged.
"""
import asyncio
import logging
import io
from datetime import datetime
from functools import partial
from typing import Optional, List
from sqlalchemy.orm import Session

from app.models import File
from app.models.ai_model import AIModel
from app.core.config import settings

logger = logging.getLogger(__name__)

# Max chars before summarization is triggered
SUMMARIZE_THRESHOLD = 80000  # ~20k tokens

# Limit concurrent extractions to avoid OOM.
# pdfplumber can use 10-50× the raw file size in working memory; allowing too
# many simultaneous extractions blows through Railway/ECS container limits.
# Keep at 1 — a single 8 MB PDF can spike to 400 MB with pdfplumber.
_EXTRACTION_SEMAPHORE = asyncio.Semaphore(1)

# Above this size, skip pdfplumber (full layout + images in RAM) and go straight
# to pypdf (streaming, ~2-3× raw file size).  3 MB is a safe threshold.
_PDFPLUMBER_MAX_BYTES = 3 * 1024 * 1024  # 3 MB

# Try to import PDF libraries (optional dependencies)
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False
    logger.warning("pdfplumber not available - PDF extraction will be limited")

try:
    from pypdf import PdfReader
    PYPDF2_AVAILABLE = True
except ImportError:
    try:
        from PyPDF2 import PdfReader
        PYPDF2_AVAILABLE = True
    except ImportError:
        PYPDF2_AVAILABLE = False
        logger.warning("pypdf/PyPDF2 not available - fallback PDF extraction disabled")

try:
    from pptx import Presentation as PptxPresentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available - PPTX extraction will use AI fallback")


class DocumentExtractionService:
    """Service for extracting text from documents and storing in database."""

    def __init__(self, db: Session):
        self.db = db

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    async def extract_and_store(
        self,
        file: File,
        force: bool = False,
        use_openai: bool = False
    ) -> bool:
        """
        Extract text from file and store in database.

        Strategy:
        1. For PDFs: try fast local extraction first (pdfplumber/PyPDF2)
        2. If local extraction fails or is poor quality, use AI extraction
        3. For non-PDF files: use AI extraction directly
        4. If content exceeds threshold, summarize with AI

        AI provider chain: DB-configured models (UseForFileExtraction=True),
        falling back to config: FILE_PROCESSING_PROVIDER → FILE_PROCESSING_FALLBACK_CHAIN

        Returns:
            True if extraction succeeded, False otherwise
        """
        # Skip if already extracted (unless force=True)
        if file.processing_status == 'ready' and not force:
            logger.debug(f"File {file.id} already processed (use force=True to re-extract)")
            return True

        logger.info(f"Extracting text from file {file.id} ({file.file_name})...")

        # Acquire semaphore — caps concurrent pdfplumber/AI work to avoid OOM.
        async with _EXTRACTION_SEMAPHORE:
            # Set processing status
            file.processing_status = "processing"
            self.db.commit()

            try:
                # Download file from blob storage
                file_bytes = await self._download_file(file)
                if not file_bytes:
                    logger.error(f"Failed to download file {file.id}")
                    file.processing_status = "failed"
                    self.db.commit()
                    return False

                is_pdf = self._is_pdf(file)
                is_pptx = self._is_pptx(file)
                extracted_text = None
                loop = asyncio.get_event_loop()

                if is_pdf:
                    # Step 1a: peek at page resources — fast, no text extraction
                    has_images = await loop.run_in_executor(
                        None, self._has_pdf_images, file_bytes
                    )

                    if has_images:
                        # PDF has embedded images/diagrams — send raw bytes to Gemini
                        # so it can read both text and visual content in one pass.
                        logger.info(f"PDF {file.id} contains images — using Gemini vision extraction")
                        extracted_text = await self._extract_pdf_with_gemini_vision(file_bytes, file)

                    if not extracted_text:
                        # Text-only PDF (or Gemini failed) — use local pypdf/pdfplumber
                        extracted_text = await loop.run_in_executor(
                            None, self._extract_with_python, file_bytes, file
                        )
                        if extracted_text and len(extracted_text.split()) < 20:
                            # Sparse result → likely a scanned PDF; try Gemini vision
                            logger.warning(
                                f"pypdf returned sparse text ({len(extracted_text.split())} words) "
                                f"for file {file.id} — trying Gemini vision"
                            )
                            extracted_text = (
                                await self._extract_pdf_with_gemini_vision(file_bytes, file)
                                or extracted_text
                            )

                elif is_pptx:
                    extracted_text = await loop.run_in_executor(
                        None, self._extract_pptx, file_bytes, file
                    )
                    if extracted_text:
                        word_count = len(extracted_text.split())
                        if word_count < 10:
                            logger.warning(f"PPTX extraction too sparse ({word_count} words), falling back to AI")
                            extracted_text = None

                # Step 2: AI extraction if local failed or non-PDF/PPTX file type
                if not extracted_text:
                    extracted_text = await self._extract_with_ai(file_bytes, file)

                # Free the raw bytes — no longer needed, can be large
                del file_bytes
                import gc; gc.collect()

                if not extracted_text:
                    logger.warning(f"No text extracted from file {file.id}")
                    file.processing_status = "failed"
                    self.db.commit()
                    return False

                # Save full extracted text to database
                word_count = len(extracted_text.split())
                file.extracted_text = extracted_text
                file.extracted_at = datetime.utcnow()
                file.extracted_word_count = word_count

                # Step 3: Summarize if content is too long, store separately
                if len(extracted_text) > SUMMARIZE_THRESHOLD:
                    logger.info(f"Content is {len(extracted_text)} chars (threshold: {SUMMARIZE_THRESHOLD}), summarizing for context use...")
                    summarized = await self._summarize_with_ai(extracted_text, file)
                    if summarized:
                        file.summarized_text = summarized
                        logger.info(f"Stored summary ({len(summarized)} chars) alongside full text ({len(extracted_text)} chars)")

                file.processing_status = "ready"
                self.db.commit()

                logger.info(f"Extracted {word_count} words from file {file.id}")
                return True

            except Exception as e:
                logger.error(f"Error extracting text from file {file.id}: {e}")
                try:
                    self.db.rollback()
                    file.processing_status = "failed"
                    self.db.commit()
                except Exception:
                    self.db.rollback()
                return False

    async def extract_all_pending(self, limit: int = 50) -> int:
        """
        Background worker: Extract text from all files without extracted_text.

        Args:
            limit: Max files to process in one batch

        Returns:
            Number of files successfully extracted
        """
        logger.info(f"Checking for files needing text extraction (limit={limit})...")

        from sqlalchemy import or_

        # Get ALL active files that haven't been successfully processed yet (not just PDFs)
        pending_files = (
            self.db.query(File)
            .filter(
                File.is_active == True,
                or_(File.processing_status != 'ready', File.processing_status == None),
                # Exclude YouTube transcripts (they have their own pipeline)
                or_(File.source_type == None, File.source_type != 'youtube')
            )
            .limit(limit)
            .all()
        )

        if not pending_files:
            logger.info("No files pending extraction")
            return 0

        logger.info(f"Found {len(pending_files)} files to extract")

        success_count = 0
        for file in pending_files:
            if await self.extract_and_store(file):
                success_count += 1

        logger.info(f"Successfully extracted {success_count}/{len(pending_files)} files")
        return success_count

    # ------------------------------------------------------------------
    # AI Provider Chain (DB-configured + config fallback)
    # ------------------------------------------------------------------

    def _get_provider_chain(self) -> List[dict]:
        """
        Get the ordered list of AI providers/models for file processing.

        Priority:
        1. DB: AIModels with UseForFileExtraction=True, ordered by provider
        2. Config fallback: FILE_PROCESSING_PROVIDER + FILE_PROCESSING_FALLBACK_CHAIN

        Returns list of dicts: [{"provider": "deepseek", "model": "deepseek-chat"}, ...]
        """
        chain = []

        # Try DB-configured models first
        try:
            db_models = (
                self.db.query(AIModel)
                .filter(
                    AIModel.use_for_file_extraction == True,
                    AIModel.is_active == True,
                    AIModel.is_deprecated == False,
                )
                .all()
            )
            if db_models:
                for model in db_models:
                    chain.append({
                        "provider": model.provider.lower(),
                        "model": model.model_name,
                    })
                logger.info(f"Using DB-configured extraction models: {[m['model'] for m in chain]}")
                return chain
        except Exception as e:
            logger.warning(f"Failed to query DB for extraction models: {e}")

        # Fallback to config settings
        primary = settings.FILE_PROCESSING_PROVIDER.lower()
        chain_str = getattr(settings, 'FILE_PROCESSING_FALLBACK_CHAIN', 'gemini,openai')
        fallbacks = [p.strip().lower() for p in chain_str.split(",") if p.strip()]

        for provider in [primary] + fallbacks:
            chain.append({"provider": provider, "model": None})

        logger.info(f"Using config-based extraction chain: {[m['provider'] for m in chain]}")
        return chain

    async def _extract_with_ai(self, file_bytes: bytes, file: File) -> Optional[str]:
        """
        Extract text using AI providers with fallback chain.

        Sends raw text content to the AI and asks it to extract/clean the content.
        For PDFs, first extracts raw text locally, then asks AI to clean it up.
        For other file types, sends the content directly.
        """
        # Get raw content to send to AI.
        # _get_raw_content may call pdfplumber/PyPDF2 (CPU-bound) so run in executor.
        loop = asyncio.get_event_loop()
        raw_content = await loop.run_in_executor(
            None, self._get_raw_content, file_bytes, file
        )
        if not raw_content:
            logger.error(f"Could not get any raw content from file {file.id}")
            return None

        # Truncate very large content to avoid token limits
        max_input = 100000  # ~25k tokens
        if len(raw_content) > max_input:
            raw_content = raw_content[:max_input] + "\n\n[... content truncated due to length ...]"

        file_type = file.content_type or file.file_type or "unknown"
        file_name = file.file_name or file.name or "document"

        extraction_prompt = f"""Extract and clean the text content from this {file_type} file named "{file_name}".

Instructions:
- Extract ALL text content, preserving structure (headings, paragraphs, lists)
- Format tables as markdown tables
- Remove any garbled characters or encoding artifacts
- Preserve the logical reading order
- Do NOT add commentary or explanations - output only the extracted text

Raw content:
---
{raw_content}
---"""

        messages = [
            {"role": "system", "content": "You are a document text extraction assistant. Extract clean, structured text from the provided content. Output ONLY the extracted text."},
            {"role": "user", "content": extraction_prompt}
        ]

        chain = self._get_provider_chain()
        provider_names = [f"{m['model'] or m['provider']}" for m in chain]
        logger.info(f"AI extraction chain for file {file.id}: {' → '.join(provider_names)}")

        for entry in chain:
            provider = entry["provider"]
            model_name = entry.get("model")
            try:
                result = await self._call_ai_provider(provider, messages, model_name=model_name)
                if result and len(result) > 50:
                    logger.info(f"AI extraction succeeded with {model_name or provider}: {len(result)} chars")
                    return result
                else:
                    logger.warning(f"AI extraction from {model_name or provider} returned insufficient content")
            except Exception as e:
                logger.warning(f"AI extraction failed with {model_name or provider}: {type(e).__name__}: {str(e)[:200]}")
                continue

        logger.error(f"All AI providers failed for file {file.id}")
        return None

    async def _summarize_with_ai(self, text: str, file: File) -> Optional[str]:
        """
        Summarize long content using AI to fit within context limits.

        Preserves key information, formulas, definitions, and examples.
        """
        file_name = file.file_name or file.name or "document"
        word_count = len(text.split())

        messages = [
            {"role": "system", "content": "You are a document summarization assistant for an educational platform. Create a comprehensive summary that preserves all key concepts, definitions, formulas, examples, and important details. The summary should be useful for answering student questions about this material."},
            {"role": "user", "content": f"""Summarize this {word_count}-word document "{file_name}" into a concise but comprehensive version.

IMPORTANT:
- Preserve ALL key concepts, definitions, theorems, and formulas
- Keep important examples and explanations
- Maintain section structure with clear headings
- Target approximately {settings.PDF_MAX_CONTEXT_CHARS // 4} words (quarter of original if possible)
- Do NOT add your own commentary

Document:
---
{text[:SUMMARIZE_THRESHOLD]}
---"""}
        ]

        chain = self._get_provider_chain()
        logger.info(f"Summarizing {word_count} words for file {file.id}")

        for entry in chain:
            provider = entry["provider"]
            model_name = entry.get("model")
            try:
                result = await self._call_ai_provider(provider, messages, model_name=model_name)
                if result and len(result) > 100:
                    summary_words = len(result.split())
                    logger.info(f"Summarized {word_count} → {summary_words} words using {model_name or provider}")
                    return result
            except Exception as e:
                logger.warning(f"Summarization failed with {model_name or provider}: {type(e).__name__}: {str(e)[:200]}")
                continue

        logger.warning("All providers failed for summarization, using raw text (truncated)")
        return None

    async def _call_ai_provider(self, provider: str, messages: list, model_name: str = None) -> Optional[str]:
        """
        Call a specific AI provider for extraction/summarization.

        Args:
            provider: Provider name (deepseek, gemini, openai, xai, etc.)
            messages: Chat messages to send
            model_name: Specific model to use (from DB), or None for default

        Each provider service has its own retry logic with exponential backoff.
        Services resolve API keys via DB KeyManager automatically.
        """
        # Create a minimal module-like object for service initialization
        from types import SimpleNamespace
        dummy_module = SimpleNamespace(
            id=0, ai_model=None, ai_model_id=None, name="extraction"
        )

        if provider == "deepseek":
            from app.services.deepseek_service import DeepSeekService
            service = DeepSeekService(module=dummy_module, model_name=model_name, db=self.db)
        elif provider == "gemini":
            from app.services.gemini_service import GeminiService
            service = GeminiService(module=dummy_module, model_name=model_name, db=self.db)
        elif provider == "xai":
            from app.services.xai_service import XAIService
            service = XAIService(module=dummy_module, model_name=model_name, db=self.db)
        elif provider == "openai":
            from app.services.ai_service import AIService
            service = AIService(module=dummy_module, db=self.db)
            if model_name:
                service.model_name = model_name
        elif provider == "anthropic":
            try:
                from app.services.anthropic_service import AnthropicService
            except ImportError:
                logger.warning("anthropic package not installed; skipping anthropic extraction provider")
                return None
            service = AnthropicService(module=dummy_module, model_name=model_name, db=self.db)
        else:
            logger.warning(f"Unknown extraction provider: {provider} — skipping")
            return None

        result = await service.answer_question_with_history(messages)
        return result.get("response", "")

    # ------------------------------------------------------------------
    # Helper Methods
    # ------------------------------------------------------------------

    def _is_pdf(self, file: File) -> bool:
        """Check if a file is a PDF."""
        return (
            (file.file_type and file.file_type.lower() in ('pdf', 'application/pdf'))
            or (file.content_type and file.content_type.lower() == 'application/pdf')
            or (file.file_name and file.file_name.lower().endswith('.pdf'))
        )

    def _is_pptx(self, file: File) -> bool:
        """Check if a file is a PowerPoint presentation."""
        pptx_types = (
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'application/vnd.ms-powerpoint',
        )
        return (
            (file.content_type and file.content_type.lower() in pptx_types)
            or (file.file_type and file.file_type.lower() in pptx_types)
            or (file.file_name and file.file_name.lower().endswith(('.pptx', '.ppt')))
        )

    def _extract_pptx(self, pptx_bytes: bytes, file: File) -> Optional[str]:
        """Extract text from a PowerPoint file using python-pptx."""
        if not PPTX_AVAILABLE:
            return None
        try:
            prs = PptxPresentation(io.BytesIO(pptx_bytes))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_texts.append(text)
                if slide_texts:
                    parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
            if parts:
                result = "\n\n".join(parts)
                logger.info(f"Extracted {len(prs.slides)} slides from PPTX file {file.id}: {len(result)} chars")
                return result
        except Exception as e:
            logger.warning(f"python-pptx extraction failed for file {file.id}: {e}")
        return None

    def _get_raw_content(self, file_bytes: bytes, file: File) -> Optional[str]:
        """
        Get raw text content from file bytes.
        Uses local extraction for PDFs/PPTX, direct decode for text files.
        """
        # Try text decode for text-based files
        if file.content_type and 'text' in file.content_type.lower():
            try:
                return file_bytes.decode('utf-8', errors='ignore')
            except Exception:
                pass

        # Try local PDF extraction
        if self._is_pdf(file):
            text = self._extract_with_python(file_bytes, file)
            if text:
                return text

        # Try PPTX extraction
        if self._is_pptx(file):
            text = self._extract_pptx(file_bytes, file)
            if text:
                return text

        # Last resort: try raw decode (handles .txt, .csv, .md, etc.)
        try:
            decoded = file_bytes.decode('utf-8', errors='ignore')
            if decoded and len(decoded.strip()) > 20:
                return decoded
        except Exception:
            pass

        # Binary files — cannot extract raw content; AI extraction won't help either
        logger.warning(f"File {file.id} ({file.file_name}) is binary with no supported extractor")
        return None

    async def _download_file(self, file: File) -> Optional[bytes]:
        """Download file from S3 using SDK, with HTTP URL fallback.

        Both boto3 and requests are synchronous libraries — we offload them to a
        thread pool via run_in_executor so they never block the event loop.
        """
        loop = asyncio.get_event_loop()

        # Primary: download via S3 SDK using blob_path (works for private buckets)
        if file.blob_path and settings.S3_BUCKET_NAME:
            try:
                import boto3

                def _s3_download() -> bytes:
                    s3_kwargs = {"region_name": settings.AWS_REGION}
                    if settings.AWS_ACCESS_KEY_ID and settings.AWS_SECRET_ACCESS_KEY:
                        s3_kwargs["aws_access_key_id"] = settings.AWS_ACCESS_KEY_ID
                        s3_kwargs["aws_secret_access_key"] = settings.AWS_SECRET_ACCESS_KEY
                    s3_client = boto3.client("s3", **s3_kwargs)
                    resp = s3_client.get_object(
                        Bucket=settings.S3_BUCKET_NAME,
                        Key=file.blob_path,
                    )
                    return resp["Body"].read()

                data = await loop.run_in_executor(None, _s3_download)
                logger.info(f"Downloaded file {file.id} from S3: {file.blob_path} ({len(data)} bytes)")
                return data
            except Exception as e:
                logger.warning(f"S3 download failed for file {file.id}: {e}, trying URL fallback...")

        # Fallback: HTTP GET on blob_url (for legacy Azure URLs or public S3 URLs)
        if file.blob_url:
            try:
                import httpx
                async with httpx.AsyncClient(timeout=30) as client:
                    response = await client.get(file.blob_url)
                    response.raise_for_status()
                    logger.info(f"Downloaded file {file.id} via URL: {len(response.content)} bytes")
                    return response.content
            except ImportError:
                # httpx not installed — fall back to blocking requests in executor
                import requests as req

                def _http_download() -> bytes:
                    r = req.get(file.blob_url, timeout=30)
                    r.raise_for_status()
                    return r.content

                data = await loop.run_in_executor(None, _http_download)
                logger.info(f"Downloaded file {file.id} via URL (executor): {len(data)} bytes")
                return data

        logger.error(f"File {file.id} has no blob_path or blob_url")
        return None

    def _has_pdf_images(self, pdf_bytes: bytes) -> bool:
        """
        Peek at PDF page resources to detect embedded images (XObjects of subtype Image).
        Uses pypdf's lazy page parser — does NOT extract text, so it's very fast and lean.
        Returns True if any page contains at least one image.
        Falls back to True on error so we don't silently skip Gemini vision.
        """
        if not PYPDF2_AVAILABLE:
            return True  # can't check — assume images so Gemini is used
        try:
            reader = PdfReader(io.BytesIO(pdf_bytes))
            for page in reader.pages:
                if page.images:
                    return True
            return False
        except Exception as e:
            logger.warning(f"Could not scan PDF for images: {e} — assuming images present")
            return True

    async def _extract_pdf_with_gemini_vision(self, pdf_bytes: bytes, file: File) -> Optional[str]:
        """
        Send the raw PDF bytes to Gemini via the native REST API so it can extract
        both text AND visual content (images, diagrams, charts, scanned pages).

        Uses inline_data (base64) which supports PDFs up to 20 MB.
        Falls back gracefully if no API key or the file is too large.
        """
        import base64
        import httpx

        # 20 MB inline_data limit from Gemini
        max_size = 20 * 1024 * 1024
        if len(pdf_bytes) > max_size:
            logger.warning(
                f"PDF {file.id} is {len(pdf_bytes) // (1024*1024)} MB — "
                f"exceeds Gemini inline limit ({max_size // (1024*1024)} MB), skipping vision"
            )
            return None

        # Resolve Gemini API key (DB-managed key first, then env)
        api_key = None
        if self.db:
            try:
                from app.services.key_manager import KeyManager
                km = KeyManager(self.db)
                result = km.get_next_key("gemini")
                if result:
                    _, api_key = result
            except Exception:
                pass
        if not api_key:
            api_key = settings.GEMINI_API_KEY
        if not api_key:
            logger.warning("No Gemini API key available — skipping PDF vision extraction")
            return None

        file_name = file.file_name or file.name or "document.pdf"
        file_size_mb = len(pdf_bytes) / (1024 * 1024)

        payload = {
            "contents": [{
                "parts": [
                    {
                        "inline_data": {
                            "mime_type": "application/pdf",
                            "data": base64.b64encode(pdf_bytes).decode(),
                        }
                    },
                    {
                        "text": (
                            f'Extract all content from this PDF document named "{file_name}".\n\n'
                            "Instructions:\n"
                            "- Extract ALL text, including text inside images, diagrams, and charts\n"
                            "- Preserve document structure (headings, paragraphs, lists)\n"
                            "- Format tables as markdown tables\n"
                            "- For diagrams or figures, describe the key information they convey\n"
                            "- Output only the extracted content — no commentary or preamble"
                        )
                    },
                ]
            }],
            "generationConfig": {
                "maxOutputTokens": 8192,
                "temperature": 0.1,
            },
        }

        # Respect DB-configured extraction model, fall back to gemini-flash-lite-latest
        model = "gemini-flash-lite-latest"
        chain = self._get_provider_chain()
        for entry in chain:
            if entry.get("provider") == "gemini" and entry.get("model"):
                model = entry["model"]
                break

        url = (
            f"https://generativelanguage.googleapis.com/v1beta"
            f"/models/{model}:generateContent?key={api_key}"
        )

        try:
            async with httpx.AsyncClient(timeout=120) as client:
                response = await client.post(url, json=payload)
                response.raise_for_status()
                data = response.json()
                text = data["candidates"][0]["content"]["parts"][0]["text"]

            if text and len(text) > 50:
                logger.info(
                    f"Gemini vision extracted {len(text)} chars from PDF {file.id} "
                    f"({file_size_mb:.1f} MB)"
                )
                return text

            logger.warning(f"Gemini vision returned insufficient content for file {file.id}")
            return None

        except Exception as e:
            logger.warning(f"Gemini vision extraction failed for file {file.id}: {type(e).__name__}: {str(e)[:200]}")
            return None

    def _extract_with_python(self, pdf_bytes: bytes, file: File) -> Optional[str]:
        """
        Extract text using Python libraries (pdfplumber or pypdf).
        Fast, free, but may struggle with complex layouts.

        pdfplumber loads full page layout (including embedded images/fonts) and
        can use 10-50× the raw file size.  For files above _PDFPLUMBER_MAX_BYTES
        we skip it entirely and use pypdf, which streams pages and stays much leaner.
        """
        import gc

        file_size = len(pdf_bytes)
        use_pdfplumber = PDFPLUMBER_AVAILABLE and file_size <= _PDFPLUMBER_MAX_BYTES

        if use_pdfplumber:
            try:
                with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
                    text_parts = []
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(page_text)

                    if text_parts:
                        result = "\n\n".join(text_parts)
                        logger.info(f"Extracted text from {len(text_parts)} pages using pdfplumber ({file_size // 1024} KB)")
                        return result
            except Exception as e:
                logger.warning(f"pdfplumber extraction failed: {e}, trying pypdf...")
            finally:
                gc.collect()
        elif PDFPLUMBER_AVAILABLE:
            logger.info(
                f"File {file.id} is {file_size // 1024} KB — skipping pdfplumber "
                f"(limit {_PDFPLUMBER_MAX_BYTES // 1024} KB), using pypdf instead"
            )

        # pypdf: preferred for large PDFs — pages are parsed lazily, ~2-3× raw size in RAM
        if PYPDF2_AVAILABLE:
            try:
                reader = PdfReader(io.BytesIO(pdf_bytes))
                text_parts = []
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)

                if text_parts:
                    logger.info(f"Extracted text from {len(text_parts)} pages using pypdf ({file_size // 1024} KB)")
                    return "\n\n".join(text_parts)
            except Exception as e:
                logger.error(f"pypdf extraction failed: {e}")
            finally:
                gc.collect()

        return None


# Convenience function
async def extract_file_text(file_id: int, db: Session, force: bool = False) -> bool:
    """Extract text from a specific file by ID."""
    file = db.query(File).filter(File.id == file_id).first()
    if not file:
        logger.error(f"File {file_id} not found")
        return False

    service = DocumentExtractionService(db)
    return await service.extract_and_store(file, force=force)
